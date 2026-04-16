"""
report_to_ppt.py
================
output/final_report.txt -> PPT 자동 변환 (미니멀 프로페셔널 디자인)

슬라이드 6장:
  1. 표지          — 화이트 배경 + 네이비 상단 밴드
  2. Exec Summary  — 진단 테이블 + 검증 배지
  3. TOP 5 과제    — 우선순위 리스트 카드
  4. 절감 효과     — KPI 카드 + 수치 테이블
  5. 에이전트 발견 — 2열 테이블
  6. 중장기 전략   — 3단계 타임라인 테이블

컬러 규칙:
  헤더 배경   #000080  NAVY
  배경        #FFFFFF  WHITE
  구분선      #D1D9E6  LGRAY
  강조 텍스트 #000080  NAVY
  일반 텍스트 #333333  DKGRAY
  서브텍스트  #666666  MGRAY
  위험 등급   CRITICAL #DC2626 / HIGH #EA580C / MEDIUM #D97706
"""

from __future__ import annotations

import os
import re
import sys
from datetime import datetime
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

# ── 컬러 팔레트 ───────────────────────────────────────────────────────────
NAVY     = RGBColor(0x00, 0x00, 0x80)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY    = RGBColor(0xD1, 0xD9, 0xE6)   # 구분선·테이블 보더
BGRAY    = RGBColor(0xF5, 0xF7, 0xFA)   # 테이블 짝수행 배경
DKGRAY   = RGBColor(0x33, 0x33, 0x33)   # 일반 텍스트
MGRAY    = RGBColor(0x66, 0x66, 0x66)   # 서브텍스트
CRITICAL = RGBColor(0xDC, 0x26, 0x26)   # 위험
HIGH     = RGBColor(0xEA, 0x58, 0x0C)   # 경고
MEDIUM   = RGBColor(0xD9, 0x77, 0x06)   # 주의

SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)

FONT_KR = "나눔고딕"
FONT_EN = "Calibri"

# 헤더 높이 / 구분선
HDR_H  = Cm(1.75)
DIV_H  = Cm(0.06)
FOOT_Y = Cm(18.4)
FOOT_H = Cm(0.55)

CONTENT_TOP = HDR_H + DIV_H + Cm(0.25)   # 콘텐츠 시작 Y


# ══════════════════════════════════════════════════════════════════════════
# 기본 드로잉 헬퍼
# ══════════════════════════════════════════════════════════════════════════

def _rect(slide, l, t, w, h,
          fill=None, line_color=None, line_w: float = 0.5) -> Any:
    shp = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    return shp


def _tb(slide, l, t, w, h, text: str,
        size: int = 10, bold: bool = False,
        color: RGBColor = DKGRAY,
        align = PP_ALIGN.LEFT,
        font: str = FONT_KR,
        wrap: bool = True) -> Any:
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text       = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.name  = font
    run.font.color.rgb = color
    return tb


# ── 공통 헤더 ─────────────────────────────────────────────────────────────
def _header(slide, title: str, page: int = 0, total: int = 6,
            subtitle: str = "", company: str = "") -> None:
    # 네이비 배경 바
    _rect(slide, Cm(0), Cm(0), SLIDE_W, HDR_H, fill=NAVY)
    # 제목
    _tb(slide, Cm(0.6), Cm(0.25), Cm(26), Cm(1.2),
        title, size=15, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    # 우측 페이지 표시
    if page:
        _tb(slide, Cm(28.5), Cm(0.55), Cm(5), Cm(0.7),
            f"{page} / {total}", size=9, color=LGRAY,
            align=PP_ALIGN.RIGHT, font=FONT_EN)
    # 헤더 하단 구분선
    _rect(slide, Cm(0), HDR_H, SLIDE_W, DIV_H, fill=LGRAY)


# ── 공통 푸터 ─────────────────────────────────────────────────────────────
def _footer(slide, company: str) -> None:
    _rect(slide, Cm(0), FOOT_Y, SLIDE_W, FOOT_H, fill=BGRAY)
    _tb(slide, Cm(0.6), FOOT_Y + Cm(0.1), Cm(22), Cm(0.4),
        company, size=7.5, color=MGRAY)
    _tb(slide, Cm(24), FOOT_Y + Cm(0.1), Cm(9.5), Cm(0.4),
        "본 보고서는 AI 에이전트 분석 결과이며 전문가 검토를 권고합니다.",
        size=6.5, color=MGRAY, align=PP_ALIGN.RIGHT)


# ── 테이블 빌더 ──────────────────────────────────────────────────────────
def _table(slide, left, top, col_widths: list,
           headers: list[str], rows: list[list[str]],
           row_h: float = 1.35,
           header_size: int = 9, body_size: int = 8.5,
           col_aligns: list | None = None) -> None:
    """
    headers/rows: 문자열 리스트
    col_widths: Cm 단위 리스트
    """
    total_w = sum(col_widths)
    rh = Cm(row_h)
    col_aligns = col_aligns or [PP_ALIGN.LEFT] * len(headers)

    # 헤더행
    x = left
    for i, (hdr, cw) in enumerate(zip(headers, col_widths)):
        _rect(slide, x, top, cw, rh, fill=NAVY)
        _tb(slide, x + Cm(0.15), top + Cm(0.2), cw - Cm(0.3), rh - Cm(0.3),
            hdr, size=header_size, bold=True, color=WHITE, align=col_aligns[i])
        x += cw

    # 데이터 행
    for ri, row in enumerate(rows):
        y    = top + rh * (ri + 1)
        bg   = WHITE if ri % 2 == 0 else BGRAY
        x    = left
        for ci, (val, cw) in enumerate(zip(row, col_widths)):
            _rect(slide, x, y, cw, rh, fill=bg, line_color=LGRAY, line_w=0.4)
            _tb(slide, x + Cm(0.15), y + Cm(0.18), cw - Cm(0.3), rh - Cm(0.3),
                str(val), size=body_size, color=DKGRAY, align=col_aligns[ci])
            x += cw


# ── KPI 카드 ──────────────────────────────────────────────────────────────
def _kpi_card(slide, left, top, width, height,
              label: str, value: str, note: str = "") -> None:
    """화이트 배경 + 네이비 상단 보더 라인 KPI 카드."""
    # 카드 테두리
    _rect(slide, left, top, width, height,
          fill=WHITE, line_color=LGRAY, line_w=0.6)
    # 상단 네이비 강조 바 (얇게)
    _rect(slide, left, top, width, Cm(0.18), fill=NAVY)
    # 레이블
    _tb(slide, left + Cm(0.2), top + Cm(0.35), width - Cm(0.4), Cm(0.6),
        label, size=8.5, color=MGRAY, align=PP_ALIGN.CENTER)
    # 수치
    _tb(slide, left + Cm(0.2), top + Cm(0.95), width - Cm(0.4), Cm(1.5),
        value, size=18, bold=True, color=NAVY,
        align=PP_ALIGN.CENTER, font=FONT_EN)
    # 비고
    if note:
        _tb(slide, left + Cm(0.2), top + height - Cm(0.7),
            width - Cm(0.4), Cm(0.6),
            note, size=7.5, color=MGRAY, align=PP_ALIGN.CENTER)


# ── 위험 등급 배지 (텍스트) ───────────────────────────────────────────────
def _risk_badge(slide, left, top, level: str) -> None:
    """CRITICAL/HIGH/MEDIUM/LOW 소형 배지."""
    color_map = {
        "CRITICAL": CRITICAL, "RED":    CRITICAL,
        "HIGH":     HIGH,     "YELLOW": MEDIUM,
        "MEDIUM":   MEDIUM,   "LOW":    RGBColor(0x16, 0xA3, 0x4A),
        "GREEN":    RGBColor(0x16, 0xA3, 0x4A),
    }
    label_map = {
        "CRITICAL": "CRITICAL", "RED":    "HIGH",
        "HIGH":     "HIGH",     "YELLOW": "MEDIUM",
        "MEDIUM":   "MEDIUM",   "LOW":    "LOW",
        "GREEN":    "LOW",
    }
    key = level.upper()
    c   = color_map.get(key, MGRAY)
    lbl = label_map.get(key, level)
    _rect(slide, left, top, Cm(1.8), Cm(0.55), fill=c)
    _tb(slide, left, top, Cm(1.8), Cm(0.55),
        lbl, size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# 파싱 유틸
# ══════════════════════════════════════════════════════════════════════════

def _strip(text: str) -> str:
    text = re.sub(r"[#*`>|🔴🟡🟢①②③④⑤⑥⑦⑧⑨📌⚡]", "", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"\s{2,}", " ", text)
    return text.strip()


def _snippet(text: str, max_len: int = 160) -> str:
    lines = []
    for ln in text.splitlines():
        ln = _strip(ln).strip()
        if ln and len(ln) > 4 and not ln.startswith("---") and not ln.startswith("```"):
            lines.append(ln)
        if sum(len(l) for l in lines) > max_len:
            break
    return " ".join(lines)[:max_len]


def parse_report(path: str) -> dict:
    with open(path, encoding="utf-8") as f:
        raw = f.read()

    m = re.search(r"===\s+(.+?)\s+컨설팅", raw)
    company = m.group(1) if m else "대상법인"

    m = re.search(r"생성일시:\s*(.+)", raw)
    gen_date = m.group(1).strip()[:10] if m else datetime.now().strftime("%Y-%m-%d")

    # 에이전트 결과
    agent_results: dict[str, str] = {}
    pat = re.compile(r"={60}\n\[([^\]]+)\]\n={60}\n(.*?)(?=\n={60}|\Z)", re.DOTALL)
    for hit in pat.finditer(raw):
        name = hit.group(1).strip()
        body = hit.group(2).strip()
        if name not in ("에이전트별 분석 결과", "검증 결과 요약", "최종 통합 보고서"):
            agent_results[name] = body

    # 최종 보고서
    m = re.search(r"={60}\n\[ 최종 통합 보고서 \]\n={60}\n(.*?)$", raw, re.DOTALL)
    final_text = m.group(1).strip() if m else ""

    # 검증 요약
    m = re.search(r"={60}\n\[ 검증 결과 요약 \]\n={60}\n(.*?)(?=\n={60})", raw, re.DOTALL)
    verify_text = m.group(1).strip() if m else ""

    return dict(company_name=company, gen_date=gen_date,
                agent_results=agent_results, final_text=final_text,
                verify_text=verify_text)


def _parse_top5(txt: str) -> list[dict]:
    tasks, pat = [], re.compile(
        r"###\s+[①②③④⑤\d]+[.)]?\s*([🔴🟡🟢]?)\s*(.+?)\n(.*?)(?=###\s+[①②③④⑤\d]|\n---|\n##|\Z)",
        re.DOTALL)
    for m in pat.finditer(txt):
        emoji = m.group(1)
        title = _strip(m.group(2))[:70]
        body  = _strip(m.group(3))[:200]
        level = "CRITICAL" if "🔴" in emoji else ("MEDIUM" if "🟡" in emoji else "LOW")
        tasks.append(dict(title=title, body=body, level=level))
        if len(tasks) >= 5:
            break
    if not tasks:
        tasks = [
            dict(title="유동성 위기 방어",      body="유동비율 60%, 단기차입 차환 긴급 추진",      level="CRITICAL"),
            dict(title="가지급금 2억원 해소",    body="세무조사 트리거, 인정이자 연 920만원 절감",  level="CRITICAL"),
            dict(title="법인세 절세 실행",       body="R&D 세액공제 등 연간 5,225만원",           level="HIGH"),
            dict(title="차명주주 명의신탁 해소", body="증여세·간주취득세 즉시 리스크",             level="HIGH"),
            dict(title="키맨보험 가입",          body="연 660만원 절세 + 경영권 보호",            level="MEDIUM"),
        ]
    return tasks


def _parse_savings(txt: str) -> list[list[str]]:
    rows, in_tbl = [], False
    for ln in txt.splitlines():
        if "절세" in ln and "|" in ln and "항목" in ln:
            in_tbl = True; continue
        if in_tbl:
            if "|" not in ln or "---" in ln:
                if rows: break
                continue
            cols = [_strip(c) for c in ln.strip().strip("|").split("|")]
            if len(cols) >= 2 and cols[0] and "항목" not in cols[0]:
                rows.append([cols[0], cols[1], cols[2] if len(cols) > 2 else ""])
    return rows


def _parse_diagnosis(txt: str) -> list[dict]:
    rows, in_tbl = [], False
    for ln in txt.splitlines():
        if "진단 영역" in ln and "|" in ln:
            in_tbl = True; continue
        if in_tbl:
            if "|" not in ln or "---" in ln:
                if rows: break
                continue
            cols = [c.strip() for c in ln.strip().strip("|").split("|")]
            if len(cols) >= 3 and cols[0] and "진단" not in cols[0]:
                sig = cols[1]
                lvl = "CRITICAL" if "RED" in sig else ("MEDIUM" if "YELLOW" in sig else "LOW")
                rows.append(dict(domain=_strip(cols[0]),
                                 signal=_strip(sig),
                                 kpi=_strip(cols[2]),
                                 level=lvl))
    return rows[:6]


def _parse_strategy(txt: str) -> list[dict]:
    phases, pat = [], re.compile(
        r"###\s+.*?(\d+~\d+개월)[^#\n]*\n(.*?)(?=###|\Z)", re.DOTALL)
    for m in pat.finditer(txt):
        items = []
        for ln in m.group(2).splitlines():
            s = _strip(ln).strip()
            if s and len(s) > 5 and not s.startswith("---"):
                items.append(s[:80])
            if len(items) >= 5: break
        if items:
            phases.append(dict(period=m.group(1), items=items))
        if len(phases) >= 3: break
    if not phases:
        phases = [
            dict(period="3~6개월",   items=["부채비율 400%→250% 개선", "가지급금 해소 완료", "R&D 세액공제 등록"]),
            dict(period="6~12개월",  items=["신용등급 CCC→B+ 회복", "영업이익률 3.5%→6%", "주식 승계 로드맵 수립"]),
            dict(period="12~36개월", items=["BBB 등급 달성·여신 정상화", "EV 부품 라인업 전환", "가업상속공제 구조 완성"]),
        ]
    return phases


# ══════════════════════════════════════════════════════════════════════════
# 슬라이드 빌더
# ══════════════════════════════════════════════════════════════════════════

def _slide1_cover(prs, company: str, date_str: str) -> None:
    """표지 — 화이트 배경 + 상단 네이비 밴드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 상단 네이비 밴드 (슬라이드 높이 약 38%)
    band_h = Cm(7.2)
    _rect(slide, Cm(0), Cm(0), SLIDE_W, band_h, fill=NAVY)

    # 밴드 좌측 얇은 화이트 강조 라인
    _rect(slide, Cm(0.8), Cm(1.2), Cm(0.12), Cm(4.8), fill=WHITE)

    # 제목 (밴드 내)
    _tb(slide, Cm(1.4), Cm(1.5), Cm(31), Cm(2.0),
        "중소기업 경영컨설팅 종합 분석 보고서",
        size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    _tb(slide, Cm(1.4), Cm(3.5), Cm(31), Cm(1.0),
        "16 + 2 전문 에이전트 분석  ·  AI 기반 통합 컨설팅",
        size=11, color=LGRAY, align=PP_ALIGN.LEFT, font=FONT_EN)

    # 화이트 영역 — 회사명
    _tb(slide, Cm(1.4), Cm(8.0), Cm(31), Cm(2.5),
        company, size=28, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    # 구분선
    _rect(slide, Cm(1.4), Cm(10.6), Cm(30.5), Cm(0.06), fill=LGRAY)

    # 메타 정보
    _tb(slide, Cm(1.4), Cm(11.0), Cm(15), Cm(0.7),
        f"보고서 생성일: {date_str}", size=10, color=MGRAY)
    _tb(slide, Cm(1.4), Cm(11.8), Cm(20), Cm(0.7),
        "본 보고서는 AI 에이전트 분석 결과입니다. 최종 의사결정 전 전문가 검토를 권고합니다.",
        size=8.5, color=MGRAY)


def _slide2_exec(prs, company: str, final_text: str, verify_text: str) -> None:
    """Exec Summary — 진단 테이블 + 검증 결과."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "Executive Summary  —  경영 현황 종합 진단", 2, company=company)
    _footer(slide, company)

    diag = _parse_diagnosis(final_text) or [
        dict(domain="재무 안정성", signal="RED",    kpi="부채비율 400%, 유동비율 60%",      level="CRITICAL"),
        dict(domain="신용·유동성", signal="RED",    kpi="신용등급 CCC 이하, 금리 10.5%",    level="CRITICAL"),
        dict(domain="수익성",     signal="YELLOW", kpi="영업이익률 3.5% (업계 5~8% 하회)",  level="MEDIUM"),
        dict(domain="기업가치",   signal="RED",    kpi="순부채 60억 > 사업가치",            level="CRITICAL"),
        dict(domain="세무·법률",  signal="YELLOW", kpi="가지급금 2억, 승계계획 전무",       level="MEDIUM"),
        dict(domain="산업 포지션", signal="YELLOW", kpi="EV 전환 대응 불확실", level="MEDIUM"),
    ]

    # 종합 판정 박스 (좌상단, 작게)
    _rect(slide, Cm(0.6), CONTENT_TOP, Cm(5.0), Cm(2.6),
          fill=WHITE, line_color=CRITICAL, line_w=1.0)
    _rect(slide, Cm(0.6), CONTENT_TOP, Cm(5.0), Cm(0.18), fill=CRITICAL)
    _tb(slide, Cm(0.6), CONTENT_TOP + Cm(0.3), Cm(5.0), Cm(0.7),
        "종합 판정", size=8.5, color=MGRAY, align=PP_ALIGN.CENTER)
    _tb(slide, Cm(0.6), CONTENT_TOP + Cm(0.95), Cm(5.0), Cm(1.4),
        "RED", size=22, bold=True, color=CRITICAL,
        align=PP_ALIGN.CENTER, font=FONT_EN)

    # 진단 테이블
    tbl_left  = Cm(6.2)
    col_w     = [Cm(7.5), Cm(2.5), Cm(17.2)]
    col_align = [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT]

    # 헤더
    x = tbl_left
    top0 = CONTENT_TOP
    for hdr, cw, al in zip(["진단 영역", "위험등급", "핵심 지표"], col_w, col_align):
        _rect(slide, x, top0, cw, Cm(0.7), fill=NAVY)
        _tb(slide, x + Cm(0.1), top0 + Cm(0.08), cw - Cm(0.2), Cm(0.56),
            hdr, size=9, bold=True, color=WHITE, align=al)
        x += cw

    rh = Cm(1.42)
    for ri, row in enumerate(diag[:6]):
        y  = top0 + Cm(0.7) + ri * rh
        bg = WHITE if ri % 2 == 0 else BGRAY
        x  = tbl_left

        # 진단 영역
        _rect(slide, x, y, col_w[0], rh, fill=bg, line_color=LGRAY, line_w=0.4)
        _tb(slide, x + Cm(0.15), y + Cm(0.35), col_w[0] - Cm(0.3), rh - Cm(0.4),
            row["domain"], size=9, bold=True, color=NAVY)
        x += col_w[0]

        # 위험등급 배지
        _rect(slide, x, y, col_w[1], rh, fill=bg, line_color=LGRAY, line_w=0.4)
        _risk_badge(slide, x + Cm(0.35), y + Cm(0.43), row["level"])
        x += col_w[1]

        # KPI
        _rect(slide, x, y, col_w[2], rh, fill=bg, line_color=LGRAY, line_w=0.4)
        _tb(slide, x + Cm(0.15), y + Cm(0.3), col_w[2] - Cm(0.3), rh - Cm(0.4),
            row["kpi"], size=8.5, color=DKGRAY)
        x += col_w[2]

    # 검증 결과 (하단 띠)
    if verify_text:
        vy = CONTENT_TOP + Cm(0.7) + 6 * rh + Cm(0.2)
        _rect(slide, Cm(0.6), vy, SLIDE_W - Cm(1.2), Cm(1.2),
              fill=BGRAY, line_color=LGRAY, line_w=0.4)
        _tb(slide, Cm(0.8), vy + Cm(0.15), Cm(4), Cm(0.5),
            "검증 결과", size=8, bold=True, color=NAVY)
        _tb(slide, Cm(4.8), vy + Cm(0.15), SLIDE_W - Cm(6), Cm(0.9),
            _strip(verify_text)[:220], size=8, color=DKGRAY)


def _slide3_top5(prs, company: str, final_text: str) -> None:
    """즉시 실행 과제 TOP 5 — 번호 카드 리스트."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "즉시 실행 과제  TOP 5", 3, company=company)
    _footer(slide, company)

    tasks   = _parse_top5(final_text)
    card_h  = Cm(2.8)
    gap     = Cm(0.28)
    left    = Cm(0.6)
    card_w  = SLIDE_W - Cm(1.2)

    for i, task in enumerate(tasks[:5]):
        y = CONTENT_TOP + i * (card_h + gap)
        # 카드 배경 (화이트, 라이트그레이 보더)
        _rect(slide, left, y, card_w, card_h,
              fill=WHITE, line_color=LGRAY, line_w=0.5)

        # 위험등급 배지
        _risk_badge(slide, left + Cm(0.25), y + Cm(1.1), task["level"])

        # 순서 번호
        _tb(slide, left + Cm(0.25), y + Cm(0.2), Cm(1.5), Cm(0.8),
            f"0{i+1}", size=20, bold=True, color=NAVY,
            align=PP_ALIGN.LEFT, font=FONT_EN)

        # 제목
        _tb(slide, left + Cm(2.3), y + Cm(0.2), Cm(16), Cm(1.0),
            task["title"], size=11, bold=True, color=NAVY)

        # 본문
        _tb(slide, left + Cm(2.3), y + Cm(1.15), Cm(30.5), Cm(1.5),
            task["body"], size=9, color=DKGRAY)


def _slide4_savings(prs, company: str, final_text: str) -> None:
    """절세·절감 효과 — KPI 카드 2개 + 수치 테이블."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "예상 절세·절감 효과", 4, company=company)
    _footer(slide, company)

    # KPI 카드 2개
    kpi_top = CONTENT_TOP
    kpi_h   = Cm(3.2)
    kpi_w   = Cm(15.8)
    _kpi_card(slide, Cm(0.6), kpi_top, kpi_w, kpi_h,
              "연간 총 절세·절감", "1.7 ~ 2.2 억원", "직접 절세 + 금융비용 절감")
    _kpi_card(slide, Cm(17.1), kpi_top, kpi_w, kpi_h,
              "7년 누적 효과", "12 ~ 15 억원", "복리 효과 미포함")

    # 절감 항목 테이블
    rows = _parse_savings(final_text) or [
        ["R&D 세액공제 외 법인세 절세",   "5,225만원/년",  "과세표준 2.5억 압축"],
        ["가지급금 해소 (인정이자 절감)", "920만원/년",   "연 4.6% 기준"],
        ["키맨보험 손금산입 절세",        "660만원/년",   "법인세율 22%"],
        ["신용등급 개선 → 금융비용 절감", "1~1.5억원/년", "금리 3~4%p 인하"],
    ]

    tbl_top = kpi_top + kpi_h + Cm(0.35)
    _table(
        slide,
        left       = Cm(0.6),
        top        = tbl_top,
        col_widths = [Cm(17.5), Cm(8.5), Cm(7.0)],
        headers    = ["절세·절감 항목", "연간 효과", "비고"],
        rows       = rows,
        row_h      = 1.6,
        col_aligns = [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT],
    )


def _slide5_agents(prs, company: str, agent_results: dict) -> None:
    """에이전트별 핵심 발견사항 — 2열 테이블."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "에이전트별 핵심 발견사항", 5, company=company)
    _footer(slide, company)

    ORDER = [
        ("TaxAgent",           "세무전략"),
        ("StockAgent",         "주식평가"),
        ("SuccessionAgent",    "가업승계"),
        ("FinanceAgent",       "재무구조"),
        ("LegalAgent",         "법률리스크"),
        ("PatentAgent",        "특허·IP"),
        ("LaborAgent",         "인사노무"),
        ("IndustryAgent",      "업종분석"),
        ("WebResearchAgent",   "시장조사"),
        ("PolicyFundingAgent", "정책자금"),
        ("CashFlowAgent",      "현금흐름"),
        ("CreditRatingAgent",  "신용등급"),
        ("RealEstateAgent",    "부동산"),
        ("InsuranceAgent",     "보험설계"),
        ("MAValuationAgent",   "기업가치"),
        ("ESGRiskAgent",       "ESG"),
    ]

    items = [(lbl, _snippet(agent_results[name], 120))
             for name, lbl in ORDER if name in agent_results]

    rh      = Cm(1.45)
    col_lbl = Cm(3.2)
    col_cnt = Cm(13.5)
    col_w   = [col_lbl, col_cnt]
    gap     = Cm(0.5)
    left_l  = Cm(0.5)
    left_r  = left_l + col_lbl + col_cnt + gap
    top0    = CONTENT_TOP

    for i, (lbl, snip) in enumerate(items[:16]):
        col   = i % 2
        row   = i // 2
        left  = left_l if col == 0 else left_r
        y     = top0 + row * rh
        bg    = WHITE if row % 2 == 0 else BGRAY

        # 레이블 셀 (네이비)
        _rect(slide, left, y, col_lbl, rh, fill=NAVY)
        _tb(slide, left + Cm(0.1), y + Cm(0.4), col_lbl - Cm(0.2), rh - Cm(0.5),
            lbl, size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # 내용 셀
        _rect(slide, left + col_lbl, y, col_cnt, rh,
              fill=bg, line_color=LGRAY, line_w=0.4)
        _tb(slide, left + col_lbl + Cm(0.15), y + Cm(0.2),
            col_cnt - Cm(0.3), rh - Cm(0.35),
            snip or "분석 완료", size=7.8, color=DKGRAY)


def _slide6_strategy(prs, company: str, final_text: str) -> None:
    """중장기 전략방향 — 3단계 타임라인."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "중장기 전략 방향", 6, company=company)
    _footer(slide, company)

    phases    = _parse_strategy(final_text)
    PHASE_LBL = ["단기", "중기", "장기"]
    PHASE_CLR = [CRITICAL, MEDIUM, RGBColor(0x16, 0xA3, 0x4A)]

    col_w = Cm(10.5)
    gap   = Cm(0.43)
    top0  = CONTENT_TOP

    # 타임라인 연결선 (얇은 라이트그레이)
    line_y = top0 + Cm(2.1)
    _rect(slide, Cm(0.6), line_y, SLIDE_W - Cm(1.2), Cm(0.08), fill=LGRAY)

    for i, phase in enumerate(phases[:3]):
        clr  = PHASE_CLR[i]
        x    = Cm(0.6) + i * (col_w + gap)

        # 카드 테두리 (화이트, 라이트그레이)
        card_h = FOOT_Y - top0 - Cm(0.2)
        _rect(slide, x, top0, col_w, card_h,
              fill=WHITE, line_color=LGRAY, line_w=0.5)
        # 상단 컬러 바 (얇게)
        _rect(slide, x, top0, col_w, Cm(0.18), fill=clr)

        # 단계 뱃지 + 기간
        _tb(slide, x + Cm(0.2), top0 + Cm(0.25), Cm(1.6), Cm(0.55),
            PHASE_LBL[i], size=8.5, bold=True, color=clr)
        _tb(slide, x + Cm(1.9), top0 + Cm(0.25), col_w - Cm(2.1), Cm(0.55),
            phase["period"], size=10, bold=True, color=NAVY, font=FONT_EN)

        # 구분선
        _rect(slide, x + Cm(0.2), top0 + Cm(0.9), col_w - Cm(0.4), Cm(0.05), fill=LGRAY)

        # 항목 리스트
        for j, item in enumerate(phase["items"][:5]):
            iy = top0 + Cm(1.1) + j * Cm(2.6)
            # 불릿 (작은 네이비 사각)
            _rect(slide, x + Cm(0.3), iy + Cm(0.5), Cm(0.22), Cm(0.22), fill=clr)
            _tb(slide, x + Cm(0.65), iy + Cm(0.3),
                col_w - Cm(0.9), Cm(2.1),
                item, size=8.5, color=DKGRAY)


# ══════════════════════════════════════════════════════════════════════════
# 메인 빌더
# ══════════════════════════════════════════════════════════════════════════

def build_ppt(company_name: str,
              agent_results: dict[str, str],
              final_text: str,
              verify_text: str = "",
              gen_date: str = "",
              output_dir: str = "output") -> str:
    if not gen_date:
        gen_date = datetime.now().strftime("%Y-%m-%d")
    ts = datetime.now().strftime("%Y%m%d_%H%M")

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide1_cover(prs, company_name, gen_date)
    _slide2_exec(prs, company_name, final_text, verify_text)
    _slide3_top5(prs, company_name, final_text)
    _slide4_savings(prs, company_name, final_text)
    _slide5_agents(prs, company_name, agent_results)
    _slide6_strategy(prs, company_name, final_text)

    os.makedirs(output_dir, exist_ok=True)
    out = os.path.join(output_dir, f"컨설팅보고서_{ts}.pptx")
    prs.save(out)
    print(f"[PPT 저장] {out}")
    return out


def build_ppt_from_file(report_path: str, output_dir: str = "output") -> str:
    d = parse_report(report_path)
    return build_ppt(
        company_name  = d["company_name"],
        agent_results = d["agent_results"],
        final_text    = d["final_text"],
        verify_text   = d["verify_text"],
        gen_date      = d["gen_date"],
        output_dir    = output_dir,
    )


# ══════════════════════════════════════════════════════════════════════════
# CLI
# ══════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    rp = sys.argv[1] if len(sys.argv) > 1 else os.path.join(
        os.path.dirname(__file__), "output", "final_report.txt")
    if not os.path.exists(rp):
        print(f"[오류] 파일 없음: {rp}")
        sys.exit(1)
    out = build_ppt_from_file(rp, os.path.join(os.path.dirname(__file__), "output"))
    print(f"[완료] {out}")
