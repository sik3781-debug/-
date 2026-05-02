"""
에이전트팀 구조도 PPT 자동 생성 스크립트
중기이코노미기업지원단 멀티에이전트 컨설팅 시스템
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os
from datetime import datetime

# ── 색상 팔레트 ─────────────────────────────────────────
NAVY      = RGBColor(0x1A, 0x37, 0x5E)   # 헤더 배경
GOLD      = RGBColor(0xC9, 0xA0, 0x2C)   # 강조 포인트
PHASE1_BG = RGBColor(0x1D, 0x4E, 0x89)   # Phase1 카드 배경
PHASE2_BG = RGBColor(0x27, 0x6A, 0x48)   # Phase2 검증 배경
PHASE3_BG = RGBColor(0x6C, 0x35, 0x7A)   # Phase3 지원 배경
ORCH_BG   = RGBColor(0x7B, 0x2D, 0x2D)   # Orchestrator 배경
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY= RGBColor(0xF0, 0xF4, 0xF8)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x2E)
SUBTITLE  = RGBColor(0x4A, 0x5A, 0x6A)

# ── 헬퍼 함수 ──────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None, line_width=Pt(0)):
    """채색된 직사각형 추가."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=12, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    """텍스트 박스 추가."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_agent_card(slide, left, top, width, height,
                   name_kr, role_kr, detail, bg_color):
    """에이전트 카드 (배경 + 이름 + 역할 + 상세)."""
    # 카드 배경
    card = add_rect(slide, left, top, width, height, bg_color,
                    line_rgb=GOLD, line_width=Pt(0.75))
    card.shadow.inherit = False

    # 에이전트명 (굵게)
    add_text_box(slide, name_kr,
                 left + 0.05, top + 0.04, width - 0.1, 0.22,
                 font_size=8.5, bold=True, color=GOLD,
                 align=PP_ALIGN.LEFT)

    # 역할
    add_text_box(slide, role_kr,
                 left + 0.05, top + 0.24, width - 0.1, 0.18,
                 font_size=7.5, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT)

    # 상세 업무
    add_text_box(slide, detail,
                 left + 0.05, top + 0.41, width - 0.1, height - 0.46,
                 font_size=6.2, bold=False, color=RGBColor(0xCC, 0xD6, 0xE0),
                 align=PP_ALIGN.LEFT)


def add_section_label(slide, text, left, top, width, height, bg_color):
    """섹션 라벨 박스."""
    add_rect(slide, left, top, width, height, bg_color)
    add_text_box(slide, text,
                 left + 0.05, top + 0.03, width - 0.1, height - 0.06,
                 font_size=8, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)


# ── 에이전트 데이터 ────────────────────────────────────

PHASE1_AGENTS = [
    # (에이전트명, 한국어 이름, 역할, 상세)
    ("TaxAgent",           "세무 에이전트",         "법인세 절세 전략 전문가",
     "법인세 절세 / 임원보수 최적화\n퇴직금 절세 / 가지급금 해결\n접대비·세액공제 한도 관리"),

    ("StockAgent",         "주식평가 에이전트",      "비상장주식 평가·주식이동 전문가",
     "보충적 평가(순손익·순자산)\n차명주식 해소 / 명의신탁 해지\n증여·양도세 최적화"),

    ("SuccessionAgent",    "가업승계 에이전트",      "가업승계 전문가",
     "가업상속공제(조특법 §18)\n창업자금 증여특례(§30의5)\n자녀법인 활용 단계적 승계"),

    ("FinanceAgent",       "재무구조 에이전트",      "재무구조 개선·리스크 헷지 전문가",
     "부채비율 개선 / 가지급금 해결\n자본잠식 해소 / 유동성 리스크\n특허·R&D 세액공제 활용"),

    ("LegalAgent",         "법률 에이전트",          "상법·형사·공정거래 법률 전문가",
     "이사회·주총 결의 하자 분석\n배임·횡령 형사 리스크\n특수관계인 공정거래 리스크"),

    ("LaborAgent",         "인사노무 에이전트",      "인사노무 전문가",
     "통상임금·퇴직급여 설계\n4대보험 두루누리 지원\n중대재해처벌법 대응 체크리스트"),

    ("InsuranceAgent",     "보험설계 에이전트",      "법인 보험 설계·세무 전문가",
     "임원·핵심인력 보험 손금산입\nD&O 책임보험 설계\n퇴직보험·보장성 보험 절세"),

    ("MAValuationAgent",   "M&A가치평가 에이전트",  "기업가치평가·M&A 전문가",
     "DCF·EV/EBITDA 가치평가\nM&A 구조 설계 및 세무\n매수·매각 Due Diligence"),

    ("PatentAgent",        "특허·IP 에이전트",       "R&D 세액공제·IP 금융 전문가",
     "연구인력개발비 세액공제\nIP 금융·기술가치평가\n특허 포트폴리오 절세 전략"),

    ("PolicyFundingAgent", "정책자금 에이전트",      "정책자금·정부지원사업 전문가",
     "중소기업 정책자금 매칭\n기술보증·신용보증 활용\n정부 R&D·보조금 신청 전략"),

    ("RealEstateAgent",    "부동산세무 에이전트",    "법인·개인 부동산 세무 전문가",
     "법인 부동산 취득·보유·양도세\n임대소득 절세 전략\n부동산 과다법인 주식 평가"),

    ("CashFlowAgent",      "현금흐름 에이전트",      "현금흐름 예측·유동성 관리 전문가",
     "월별 현금흐름 예측 모델\n유동성 위기 시나리오 시뮬레이션\n운전자본 최적화 전략"),

    ("CreditRatingAgent",  "신용등급 에이전트",      "신용등급 분석·금융비용 절감 전문가",
     "NICE·KCB 신용등급 분석\n금리 절감 전략 도출\n여신 한도 확대 로드맵"),

    ("IndustryAgent",      "업종분석 에이전트",      "동종업계 벤치마크·세무조사 트렌드 분석가",
     "동종업계 재무비율 벤치마크\n세무조사 빈도·트렌드 분석\n업종별 절세 포인트"),

    ("ESGRiskAgent",       "ESG리스크 에이전트",     "ESG 리스크 진단 전문가",
     "환경·사회·지배구조 리스크 진단\nESG 공시 대응 전략\n지속가능경영 보고서 초안"),

    ("WebResearchAgent",   "웹리서치 에이전트",      "기업 정보 수집 리서처",
     "기업·업종 공개 정보 수집\n실시간 뉴스·공시 모니터링\n경쟁사 동향 분석"),
]

PHASE2_AGENTS = [
    ("VerifyTax",      "세무검증 에이전트",       "세무·주식평가·재무 검증 전문가",
     "법인세·상증세 법령 정확성\n비상장주식 보충적 평가 산식\n가업상속공제 요건·한도"),

    ("VerifyOps",      "운영검증 에이전트",       "법률·노무·특허·부동산·보험 검증 전문가",
     "상법 절차·요건 정확성\n노동법·근로기준법 조문\n최신 판례·예규 반영 여부"),

    ("VerifyStrategy", "전략검증 에이전트",       "경영전략·정책자금·ESG·M&A 검증 전문가",
     "M&A 가치평가 적절성\n정책자금 요건·한도 검증\nESG 공시 기준 정합성"),
]

SUPPORT_AGENTS = [
    ("DataValidationAgent", "데이터검증 에이전트", "재무데이터 정합성 검증 전문가",
     "Phase 1 최우선 실행\n입력값 유효성·범위 검사\n이상값 감지 및 보정 권고"),

    ("TaxLawUpdateAgent",   "세법업데이트 에이전트", "세법 개정사항 수집·요약 전문가",
     "국세청·기재부 최신 개정사항\n웹 검색 기반 실시간 업데이트\nPhase 2 최우선 실행"),

    ("AutoFixAgent",        "자동수정 에이전트",   "에이전트 오류 분석 및 자동 수정 전문가",
     "VerifyAgent FAIL 항목 분석\n원인 진단 후 재실행 루프\nPhase 3 검증 직후 실행"),

    ("MonitorAgent",        "모니터링 에이전트",   "경영지표 Delta 분석 전문가",
     "이전 스냅샷↔현재 지표 비교\n개선·악화 항목 정량 보고\n다음 점검 일정 권고"),

    ("ScenarioAgent",       "시나리오 에이전트",   "절세·구조설계 시나리오 시뮬레이션 전문가",
     "임원퇴직금 시나리오 3종\n가업승계·자녀법인 수치 비교\n최적안 자동 선택 도출"),
]


# ── 슬라이드 1 : 표지 ─────────────────────────────────

def make_cover(prs):
    slide_layout = prs.slide_layouts[6]  # 완전 빈 레이아웃
    slide = prs.slides.add_slide(slide_layout)

    W, H = 13.33, 7.5

    # 전체 배경
    add_rect(slide, 0, 0, W, H, NAVY)

    # 골드 상단 띠
    add_rect(slide, 0, 0, W, 0.12, GOLD)

    # 골드 하단 띠
    add_rect(slide, 0, H - 0.12, W, 0.12, GOLD)

    # 중앙 제목 영역
    add_rect(slide, 1.2, 1.5, W - 2.4, 4.5, RGBColor(0x0D, 0x1B, 0x3E),
             line_rgb=GOLD, line_width=Pt(1.5))

    # 메인 타이틀
    add_text_box(slide, "멀티에이전트 컨설팅 시스템",
                 1.4, 2.0, W - 2.8, 1.0,
                 font_size=30, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)

    # 서브 타이틀
    add_text_box(slide, "에이전트팀 구성 및 업무 체계도",
                 1.4, 3.1, W - 2.8, 0.7,
                 font_size=20, bold=False, color=GOLD,
                 align=PP_ALIGN.CENTER)

    # 구분선 (골드)
    add_rect(slide, 3.5, 3.9, 6.3, 0.04, GOLD)

    # 기관명
    add_text_box(slide, "중기이코노미기업지원단",
                 1.4, 4.1, W - 2.8, 0.5,
                 font_size=14, bold=False, color=RGBColor(0xAA, 0xBB, 0xCC),
                 align=PP_ALIGN.CENTER)

    # 날짜
    today = datetime.today().strftime("%Y년 %m월 %d일")
    add_text_box(slide, today,
                 1.4, 4.7, W - 2.8, 0.4,
                 font_size=11, bold=False, color=RGBColor(0x88, 0x99, 0xAA),
                 align=PP_ALIGN.CENTER)

    # 에이전트 수 뱃지
    for i, (num, label, col) in enumerate([
        ("25", "총 에이전트", PHASE1_BG),
        ("16", "분석 에이전트", PHASE1_BG),
        ("3",  "검증 에이전트", PHASE2_BG),
        ("5",  "지원 에이전트", PHASE3_BG),
    ]):
        x = 1.4 + i * 2.65
        add_rect(slide, x, 5.4, 2.3, 0.9, col,
                 line_rgb=GOLD, line_width=Pt(0.5))
        add_text_box(slide, num,
                     x, 5.45, 2.3, 0.45,
                     font_size=22, bold=True, color=GOLD,
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, label,
                     x, 5.9, 2.3, 0.3,
                     font_size=8, bold=False, color=WHITE,
                     align=PP_ALIGN.CENTER)


# ── 슬라이드 2 : 전체 구조도 ──────────────────────────

def make_overview(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    W, H = 13.33, 7.5

    # 배경
    add_rect(slide, 0, 0, W, H, LIGHT_GRAY)

    # 헤더
    add_rect(slide, 0, 0, W, 0.75, NAVY)
    add_text_box(slide, "시스템 구조 개요",
                 0.2, 0.1, 9, 0.55,
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text_box(slide, "4 Phase 실행 체계 · 25개 전문 에이전트",
                 0.2, 0.42, 9, 0.3,
                 font_size=9, bold=False, color=GOLD, align=PP_ALIGN.LEFT)

    # ── Orchestrator (최상위) ──
    add_rect(slide, 4.5, 0.9, 4.33, 0.65, ORCH_BG,
             line_rgb=GOLD, line_width=Pt(1.2))
    add_text_box(slide, "🎯  Orchestrator  (run.py)",
                 4.5, 0.95, 4.33, 0.55,
                 font_size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # 화살표 (세로 줄)
    add_rect(slide, 6.615, 1.55, 0.1, 0.3, NAVY)

    # ── Phase 박스들 ──
    phases = [
        ("Phase 1\n분석 에이전트 16종\n(병렬 실행)", 0.3,  1.95, 3.3, 1.3, PHASE1_BG),
        ("Phase 2\n검증 에이전트 3종\n(병렬 검증)", 3.8,  1.95, 3.3, 1.3, PHASE2_BG),
        ("Phase 3\nReportAgent\n최종 보고서 생성",  7.3,  1.95, 2.5, 1.3, NAVY),
        ("Phase 4\nPPT 자동 변환\noutput/ 저장",   10.05, 1.95, 2.9, 1.3, ORCH_BG),
    ]
    for label, x, y, w, h, bg in phases:
        add_rect(slide, x, y, w, h, bg, line_rgb=GOLD, line_width=Pt(0.75))
        add_text_box(slide, label, x + 0.1, y + 0.1, w - 0.2, h - 0.2,
                     font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Phase 연결 화살표
    for ax in [3.6, 7.1, 9.8]:
        add_rect(slide, ax, 2.45, 0.2, 0.1, GOLD)

    # ── Phase 1 세부 에이전트 그리드 (8×2) ──
    add_rect(slide, 0.2, 3.4, 12.9, 0.3, PHASE1_BG)
    add_text_box(slide, "▶  Phase 1  분석 에이전트 (16종) — ThreadPoolExecutor 병렬 실행",
                 0.3, 3.41, 12.7, 0.28,
                 font_size=8.5, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    p1_names = [
        "TaxAgent\n법인세 절세",
        "StockAgent\n비상장주식",
        "SuccessionAgent\n가업승계",
        "FinanceAgent\n재무구조",
        "LegalAgent\n법률리스크",
        "LaborAgent\n인사노무",
        "InsuranceAgent\n법인보험",
        "MAValuationAgent\nM&A평가",
        "PatentAgent\nR&D·IP",
        "PolicyFundingAgent\n정책자금",
        "RealEstateAgent\n부동산세무",
        "CashFlowAgent\n현금흐름",
        "CreditRatingAgent\n신용등급",
        "IndustryAgent\n업종분석",
        "ESGRiskAgent\nESG리스크",
        "WebResearchAgent\n기업정보",
    ]
    cols, rows = 8, 2
    card_w = 12.9 / cols
    card_h = 0.78
    for i, nm in enumerate(p1_names):
        col = i % cols
        row = i // cols
        x = 0.2 + col * card_w
        y = 3.75 + row * (card_h + 0.05)
        add_rect(slide, x, y, card_w - 0.05, card_h, PHASE1_BG,
                 line_rgb=RGBColor(0x5A, 0x80, 0xB0), line_width=Pt(0.4))
        add_text_box(slide, nm,
                     x + 0.04, y + 0.04, card_w - 0.13, card_h - 0.08,
                     font_size=6.8, bold=False, color=WHITE,
                     align=PP_ALIGN.CENTER)

    # ── Phase 2 검증 에이전트 ──
    add_rect(slide, 0.2, 5.45, 12.9, 0.28, PHASE2_BG)
    add_text_box(slide, "▶  Phase 2  검증 에이전트 (3종) — 세무 / 운영 / 전략",
                 0.3, 5.46, 12.7, 0.26,
                 font_size=8.5, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    v_agents = [
        ("VerifyTax\n세무·주식·재무 검증", PHASE2_BG),
        ("VerifyOps\n법률·노무·특허·부동산 검증", PHASE2_BG),
        ("VerifyStrategy\n전략·정책자금·ESG·M&A 검증", PHASE2_BG),
    ]
    for i, (nm, bg) in enumerate(v_agents):
        x = 0.2 + i * 4.3
        add_rect(slide, x, 5.78, 4.2, 0.55, bg,
                 line_rgb=GOLD, line_width=Pt(0.5))
        add_text_box(slide, nm, x + 0.08, 5.82, 4.1, 0.47,
                     font_size=8, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)

    # ── 지원 에이전트 (하단 우측) ──
    s_agents = [
        "DataValidationAgent\n데이터 정합성 검증",
        "TaxLawUpdateAgent\n세법 개정 업데이트",
        "AutoFixAgent\n오류 자동 수정",
        "MonitorAgent\n경영지표 Delta 분석",
        "ScenarioAgent\n시나리오 시뮬레이션",
    ]
    support_x_start = 0.2
    support_w = 2.54
    for i, nm in enumerate(s_agents):
        x = support_x_start + i * (support_w + 0.05)
        add_rect(slide, x, 6.42, support_w, 0.6, PHASE3_BG,
                 line_rgb=GOLD, line_width=Pt(0.4))
        add_text_box(slide, nm, x + 0.06, 6.44, support_w - 0.12, 0.56,
                     font_size=6.8, bold=False, color=WHITE,
                     align=PP_ALIGN.CENTER)

    add_rect(slide, 0.2, 6.38, 12.9, 0.02, PHASE3_BG)
    add_text_box(slide, "▶  지원 에이전트 (5종) — DataValidation · TaxLawUpdate · AutoFix · Monitor · Scenario",
                 0.3, 6.34, 12.7, 0.24,
                 font_size=7.5, bold=True, color=PHASE3_BG, align=PP_ALIGN.LEFT)


# ── 슬라이드 3 : Phase 1 분석 에이전트 상세 ──────────

def make_phase1_detail(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    W, H = 13.33, 7.5

    add_rect(slide, 0, 0, W, H, LIGHT_GRAY)
    add_rect(slide, 0, 0, W, 0.75, NAVY)
    add_text_box(slide, "Phase 1 — 분석 에이전트 상세 (16종)",
                 0.2, 0.1, 10, 0.55,
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text_box(slide, "ThreadPoolExecutor · max_workers=2 배치 병렬 실행",
                 0.2, 0.42, 10, 0.3,
                 font_size=9, bold=False, color=GOLD, align=PP_ALIGN.LEFT)

    cols = 4
    card_w = (W - 0.5) / cols - 0.08
    card_h = 1.55
    margin_x = 0.2
    margin_top = 0.85

    for i, (eng, kr, role, detail) in enumerate(PHASE1_AGENTS):
        col = i % cols
        row = i // cols
        x = margin_x + col * (card_w + 0.08)
        y = margin_top + row * (card_h + 0.08)
        add_agent_card(slide, x, y, card_w, card_h,
                       f"{kr}  [{eng}]", role, detail, PHASE1_BG)


# ── 슬라이드 4 : 검증·지원 에이전트 상세 ─────────────

def make_verify_support_detail(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    W, H = 13.33, 7.5

    add_rect(slide, 0, 0, W, H, LIGHT_GRAY)
    add_rect(slide, 0, 0, W, 0.75, NAVY)
    add_text_box(slide, "Phase 2·3 — 검증 · 지원 에이전트 상세 (8종)",
                 0.2, 0.1, 10, 0.55,
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text_box(slide, "검증(VerifyTax / VerifyOps / VerifyStrategy) + 운영 지원(DataValidation · TaxLawUpdate · AutoFix · Monitor · Scenario)",
                 0.2, 0.42, 12.9, 0.3,
                 font_size=8, bold=False, color=GOLD, align=PP_ALIGN.LEFT)

    # 검증 에이전트 (3종)
    add_rect(slide, 0.2, 0.85, W - 0.4, 0.25, PHASE2_BG)
    add_text_box(slide, "◆  검증 에이전트 (Phase 2 — 병렬 검증)",
                 0.3, 0.86, W - 0.5, 0.23,
                 font_size=8.5, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    vw = (W - 0.6) / 3 - 0.08
    for i, (eng, kr, role, detail) in enumerate(PHASE2_AGENTS):
        x = 0.2 + i * (vw + 0.08)
        add_agent_card(slide, x, 1.15, vw, 1.8,
                       f"{kr}  [{eng}]", role, detail, PHASE2_BG)

    # 지원 에이전트 (5종)
    add_rect(slide, 0.2, 3.1, W - 0.4, 0.25, PHASE3_BG)
    add_text_box(slide, "◆  지원 에이전트 (Phase 별 특수 목적 실행)",
                 0.3, 3.11, W - 0.5, 0.23,
                 font_size=8.5, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    sw = (W - 0.6) / 3 - 0.08
    sh = 1.85
    positions = [
        (0, 0), (1, 0), (2, 0),
        (0, 1), (1, 1),
    ]
    for i, (eng, kr, role, detail) in enumerate(SUPPORT_AGENTS):
        col, row = positions[i]
        x = 0.2 + col * (sw + 0.08)
        y = 3.4 + row * (sh + 0.08)
        add_agent_card(slide, x, y, sw, sh,
                       f"{kr}  [{eng}]", role, detail, PHASE3_BG)


# ── 슬라이드 5 : 실행 흐름 (타임라인) ────────────────

def make_flow(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    W, H = 13.33, 7.5

    add_rect(slide, 0, 0, W, H, NAVY)
    add_rect(slide, 0, 0, W, 0.12, GOLD)
    add_rect(slide, 0, H - 0.12, W, 0.12, GOLD)

    add_text_box(slide, "실행 흐름 — 4 Phase 자동화 파이프라인",
                 0.5, 0.25, 12, 0.55,
                 font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "run.py → orchestrator.py → agents/ → output/",
                 0.5, 0.78, 12, 0.35,
                 font_size=11, bold=False, color=GOLD, align=PP_ALIGN.CENTER)

    phases = [
        ("Phase 0\n초기화",
         "DataValidationAgent\n입력값 정합성 검증\n이상값 감지·보정",
         RGBColor(0x45, 0x45, 0x75)),
        ("Phase 1\n분석",
         "16개 전문 에이전트\nThreadPoolExecutor\n병렬 분석 실행",
         PHASE1_BG),
        ("Phase 2\n검증",
         "VerifyTax · VerifyOps\nVerifyStrategy\n3종 병렬 검증",
         PHASE2_BG),
        ("Phase 3\n보고서",
         "AutoFixAgent\nReportAgent\n통합 보고서 생성",
         RGBColor(0x55, 0x35, 0x75)),
        ("Phase 4\nPPT 변환",
         "report_to_ppt_navy.js\nPPT 자동 생성\noutput/ 저장",
         ORCH_BG),
    ]

    box_w = 2.3
    gap = 0.25
    total_w = len(phases) * box_w + (len(phases) - 1) * gap
    start_x = (W - total_w) / 2

    for i, (phase_label, detail, bg) in enumerate(phases):
        x = start_x + i * (box_w + gap)
        y = 1.35

        # 연결 화살표
        if i > 0:
            ax = x - gap
            add_rect(slide, ax, y + 0.9, gap, 0.1, GOLD)

        # 박스
        add_rect(slide, x, y, box_w, 2.0, bg,
                 line_rgb=GOLD, line_width=Pt(1.0))

        # 페이즈 라벨
        add_text_box(slide, phase_label,
                     x + 0.1, y + 0.1, box_w - 0.2, 0.55,
                     font_size=11, bold=True, color=GOLD,
                     align=PP_ALIGN.CENTER)

        # 상세
        add_text_box(slide, detail,
                     x + 0.1, y + 0.7, box_w - 0.2, 1.2,
                     font_size=8.5, bold=False, color=WHITE,
                     align=PP_ALIGN.CENTER)

    # 지원 에이전트 상시 실행 박스
    add_rect(slide, start_x, 3.55, total_w, 1.3,
             RGBColor(0x3A, 0x1A, 0x5A),
             line_rgb=PHASE3_BG, line_width=Pt(1.0))
    add_text_box(slide, "상시 지원 에이전트",
                 start_x + 0.15, 3.6, 2.2, 0.4,
                 font_size=9, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    support_names = [
        "MonitorAgent\n경영지표 Delta 분석",
        "ScenarioAgent\n절세 시나리오 시뮬레이션",
        "TaxLawUpdateAgent\n세법 개정 실시간 업데이트",
        "WebResearchAgent\n기업 정보 수집 리서처",
        "AutoFixAgent\n오류 자동 수정 루프",
    ]
    sw = total_w / len(support_names)
    for i, nm in enumerate(support_names):
        x = start_x + i * sw + 0.05
        add_rect(slide, x, 4.05, sw - 0.1, 0.65,
                 PHASE3_BG,
                 line_rgb=RGBColor(0xAA, 0x80, 0xCC), line_width=Pt(0.4))
        add_text_box(slide, nm, x + 0.05, 4.08, sw - 0.2, 0.6,
                     font_size=7, bold=False, color=WHITE,
                     align=PP_ALIGN.CENTER)

    # 기술 스택
    add_rect(slide, start_x, 5.05, total_w, 1.7,
             RGBColor(0x0D, 0x1B, 0x3E),
             line_rgb=RGBColor(0x33, 0x55, 0x88), line_width=Pt(0.5))
    add_text_box(slide, "기술 스택",
                 start_x + 0.15, 5.1, 2, 0.35,
                 font_size=9, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    tech_items = [
        ("모델", "claude-sonnet-4-6\nclaude-haiku-4-5"),
        ("SDK", "anthropic Python SDK\nPrompt Caching 적용"),
        ("병렬처리", "ThreadPoolExecutor\nmax_workers=2 배치"),
        ("PPT변환", "report_to_ppt_navy.js\n(Node.js · pptxgenjs)"),
        ("출력", "output/ 디렉터리\nPPT · JSON 저장"),
    ]
    tw = total_w / len(tech_items)
    for i, (t_label, t_detail) in enumerate(tech_items):
        tx = start_x + i * tw + 0.05
        add_text_box(slide, t_label,
                     tx, 5.5, tw - 0.1, 0.3,
                     font_size=8, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text_box(slide, t_detail,
                     tx, 5.82, tw - 0.1, 0.8,
                     font_size=7.5, bold=False, color=WHITE, align=PP_ALIGN.CENTER)


# ── 메인 실행 ─────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    make_cover(prs)
    make_overview(prs)
    make_phase1_detail(prs)
    make_verify_support_detail(prs)
    make_flow(prs)

    os.makedirs("output", exist_ok=True)
    out_path = os.path.join("output", "에이전트팀_구조도.pptx")
    prs.save(out_path)
    print(f"[완료] PPT 저장: {out_path}")


if __name__ == "__main__":
    main()
