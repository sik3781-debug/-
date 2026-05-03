"""
agents/active/ppt_polisher.py
===============================
PPTPolisher — 중기이코노미 네이비 표준 검증 + 자동수정
python-pptx 없이도 동작 (구조 검사 모드 폴백)
"""
from __future__ import annotations

import json
import os
import re
from datetime import datetime
from typing import Any

_ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
_LOG  = os.path.join(_ROOT, "storage", "ppt_polisher_log.jsonl")

# 금지 문구 (CLAUDE.md 기준)
BANNED_PHRASES = [
    "2025년 귀속 세법 기준",
    "경력 약 15년",
]

# 표준 색상 (hex # 없음)
STANDARD_COLORS = {
    "main":    "000080",
    "bg":      "FFFFFF",
    "divider": "D1D9E6",
    "accent":  "1E40AF",
    "tax":     "B45309",
    "finance": "0D9488",
    "succession": "9F1239",
    "strategy":   "1E40AF",
    "verify":     "047857",
    "external":   "334155",
}

FORBIDDEN_COLOR_PATTERNS = ["FFD700", "FFC000", "F9C74F"]  # 골드 계열


class PPTPolisher:

    def validate(self, pptx_path: str, mode: str = "standard") -> dict[str, Any]:
        if not os.path.exists(pptx_path):
            return {"status": "file_not_found", "path": pptx_path}

        try:
            import pptx
            return self._validate_with_pptx(pptx_path, mode)
        except ImportError:
            return self._validate_structural(pptx_path)

    def _validate_with_pptx(self, path: str, mode: str) -> dict[str, Any]:
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation(path)
        violations: list[dict] = []
        slide_count = len(prs.slides)

        for idx, slide in enumerate(prs.slides, 1):
            shape_count = len(slide.shapes)
            has_image = any(
                hasattr(s, "image") for s in slide.shapes
            )
            text_only = shape_count <= 1 and not has_image
            if text_only and mode in ("standard", "strict"):
                violations.append({
                    "slide": idx,
                    "rule": "text_only_slide",
                    "auto_fixable": False,
                    "message": f"슬라이드 {idx}: 텍스트 단독 슬라이드 (도형 {shape_count}개)",
                })

            # 금지 문구 스캔
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    for phrase in BANNED_PHRASES:
                        if phrase in shape.text:
                            violations.append({
                                "slide": idx,
                                "rule": "banned_phrase",
                                "auto_fixable": False,
                                "message": f"금지 문구 발견: '{phrase}'",
                            })

        if not (5 <= slide_count <= 30):
            violations.append({
                "slide": "all",
                "rule": "slide_count",
                "auto_fixable": False,
                "message": f"슬라이드 {slide_count}장 — 권장 5~30장",
            })

        passed = len(violations) == 0
        result = {
            "status": "pass" if passed else "violations_found",
            "path": path,
            "slide_count": slide_count,
            "mode": mode,
            "violations": violations,
            "violation_count": len(violations),
            "validated_at": datetime.now().isoformat(timespec="seconds"),
        }
        self._save_log(result)
        return result

    def _validate_structural(self, path: str) -> dict[str, Any]:
        """python-pptx 없을 때 파일 크기·확장자 기본 검사."""
        file_size = os.path.getsize(path)
        violations = []

        if not path.lower().endswith((".pptx", ".ppt")):
            violations.append({"rule": "file_extension", "message": ".pptx/.ppt 아님"})
        if file_size < 10_000:
            violations.append({"rule": "file_size", "message": f"파일 크기 {file_size}B — 내용 없음 의심"})

        result = {
            "status": "structural_only",
            "path": path,
            "file_size_bytes": file_size,
            "violations": violations,
            "violation_count": len(violations),
            "note": "python-pptx 미설치 — 구조 검사만 수행 (pip install python-pptx)",
            "validated_at": datetime.now().isoformat(timespec="seconds"),
        }
        self._save_log(result)
        return result

    def _save_log(self, result: dict) -> None:
        try:
            os.makedirs(os.path.dirname(_LOG), exist_ok=True)
            with open(_LOG, "a", encoding="utf-8") as f:
                f.write(json.dumps(result, ensure_ascii=False) + "\n")
        except Exception:
            pass

    def analyze(self, data: dict) -> dict:
        pptx_path = data.get("pptx_path", "")
        mode = data.get("mode", "standard")
        if not pptx_path:
            return {"status": "no_path", "message": "pptx_path 필요"}
        return self.validate(pptx_path, mode)
